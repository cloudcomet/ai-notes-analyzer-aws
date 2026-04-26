import boto3
import os
import sys
import tempfile

import pdfplumber
import PyPDF2
from docx import Document
from pptx import Presentation

# -------------------------
# AWS CLIENT
# -------------------------
s3 = boto3.client("s3")


# -------------------------
# CLEAN TEXT
# -------------------------
def clean_text(text):
    if not text:
        return ""

    text = text.replace("\x00", " ")
    text = text.replace("ï¿½", "")
    text = text.replace("\n\n", "\n")

    return text.strip()


# -------------------------
# PDF PARSER
# -------------------------
def parse_pdf(file_path):
    print("Parsing PDF with pdfplumber...")

    text = ""

    try:
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"

        if text.strip():
            return clean_text(text)

        print("pdfplumber empty, fallback to PyPDF2...")

    except Exception as e:
        print("pdfplumber failed:", str(e))

    # fallback
    try:
        reader = PyPDF2.PdfReader(file_path)
        for page in reader.pages:
            text += page.extract_text() or ""

        return clean_text(text)

    except Exception as e:
        print("PyPDF2 failed:", str(e))
        return "PDF extraction failed"


# -------------------------
# DOCX PARSER
# -------------------------
def parse_docx(file_path):
    print("Parsing DOCX...")

    try:
        doc = Document(file_path)
        text = "\n".join([p.text for p in doc.paragraphs])
        return clean_text(text)

    except Exception as e:
        print("DOCX failed:", str(e))
        return "DOCX extraction failed"


# -------------------------
# PPTX PARSER
# -------------------------
def parse_pptx(file_path):
    print("Parsing PPTX...")

    try:
        prs = Presentation(file_path)
        text = ""

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

        return clean_text(text)

    except Exception as e:
        print("PPTX failed:", str(e))
        return "PPTX extraction failed"


# -------------------------
# ROUTER
# -------------------------
def parse_file(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    print("Detected extension:", ext)

    if ext == ".pdf":
        return parse_pdf(file_path)

    elif ext == ".docx":
        return parse_docx(file_path)

    elif ext == ".pptx":
        return parse_pptx(file_path)

    elif ext == ".txt":
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return clean_text(f.read())

    else:
        return "Unsupported file type"


# -------------------------
# MAIN
# -------------------------
def main():
    if len(sys.argv) < 3:
        print("Usage: process.py <bucket> <key>")
        return

    bucket = sys.argv[1]
    key = sys.argv[2]

    print("Processing file:", key)

    # Skip already processed files
    if key.startswith("processed/"):
        print("Skipping processed file")
        return

    try:
        # -------------------------
        # PRESERVE EXTENSION (CRITICAL FIX)
        # -------------------------
        ext = os.path.splitext(key)[1]
        print("Original extension:", ext)

        # -------------------------
        # DOWNLOAD FILE
        # -------------------------
        with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
            s3.download_fileobj(bucket, key, tmp)
            local_path = tmp.name

        print("Downloaded file to:", local_path)

        # -------------------------
        # PARSE FILE
        # -------------------------
        text = parse_file(local_path)

        print("----- EXTRACTED TEXT PREVIEW -----")
        print(text[:500])

        # -------------------------
        # SAVE OUTPUT
        # -------------------------
        output_key = "processed/" + os.path.basename(key) + ".txt"

        s3.put_object(
            Bucket=bucket,
            Key=output_key,
            Body=text.encode("utf-8")
        )

        print("Uploaded to:", output_key)

    except Exception as e:
        print("ERROR:", str(e))


# -------------------------
# ENTRY
# -------------------------
if __name__ == "__main__":
    main()
